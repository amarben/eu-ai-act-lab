#!/usr/bin/env python3
"""
Generate comprehensive PowerPoint presentation from EU AI Act course lectures.
Uses direct text box approach for maximum reliability.
"""

import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def extract_lecture_data(lecture_file):
    """Extract structured data from a lecture markdown file."""
    with open(lecture_file, 'r', encoding='utf-8') as f:
        content = f.read()

    # Extract lecture title
    title_match = re.search(r'^# Lecture (\d+): (.+)$', content, re.MULTILINE)
    if title_match:
        lecture_num = title_match.group(1)
        lecture_title = title_match.group(2)
    else:
        lecture_num = "?"
        lecture_title = "Untitled Lecture"

    # Extract duration
    duration_match = re.search(r'Duration: (.+)$', content, re.MULTILINE)
    duration = duration_match.group(1) if duration_match else "Unknown"

    slides = []

    # Find all slides in the lecture
    slide_sections = re.finditer(r'## Slide \d+:.*?(?=## Slide \d+:|## Practice Assignment|## Final|## Downloadable|## Key Takeaways|$)', content, re.DOTALL)

    for slide_match in slide_sections:
        slide_content = slide_match.group(0)

        # Extract slide title (with or without quotes)
        title_pattern = r'### Slide Title\s*\n["\']?([^\n"\']+)["\']?'
        title_match = re.search(title_pattern, slide_content)
        if title_match:
            slide_title = title_match.group(1).strip().strip('"\'')
        else:
            slide_title = "Untitled Slide"

        # Extract slide notes (narration)
        notes_pattern = r'### Slide Notes.*?\n\n(.+?)(?=\n###|\Z)'
        notes_match = re.search(notes_pattern, slide_content, re.DOTALL)
        narration = notes_match.group(1).strip() if notes_match else ""

        # Extract bullet points
        bullets_pattern = r'### Bullet Points\s*\n\n((?:- .+\n?)+)'
        bullets_match = re.search(bullets_pattern, slide_content)

        bullet_points = []
        if bullets_match:
            bullet_text = bullets_match.group(1)
            bullet_points = [line.strip('- ').strip() for line in bullet_text.strip().split('\n') if line.strip().startswith('-')]

        slides.append({
            'title': slide_title,
            'bullets': bullet_points,
            'narration': narration
        })

    return {
        'lecture_num': lecture_num,
        'lecture_title': lecture_title,
        'duration': duration,
        'slides': slides
    }

def create_title_slide(prs):
    """Create course title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title text box
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1.5)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "EU AI Act Compliance"

    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # Add subtitle text box
    left = Inches(1)
    top = Inches(3.8)
    width = Inches(8)
    height = Inches(2)

    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Build Audit-Ready Documentation\n\nComplete Course Presentation\n9 Lectures | 27 Slides"

    for paragraph in subtitle_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
        paragraph.alignment = PP_ALIGN.CENTER

def create_section_divider(prs, lecture_num, lecture_title, duration):
    """Create section divider slide between lectures."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title text box
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = f"Lecture {lecture_num}: {lecture_title}"

    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # Add duration text box
    left = Inches(1)
    top = Inches(4.2)
    width = Inches(8)
    height = Inches(0.8)

    duration_box = slide.shapes.add_textbox(left, top, width, height)
    duration_frame = duration_box.text_frame
    duration_frame.text = f"Duration: {duration}"

    p = duration_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(96, 96, 96)
    p.alignment = PP_ALIGN.CENTER

def create_content_slide(prs, lecture_num, slide_data):
    """Create content slide with title, bullets, and speaker notes using text boxes."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title text box
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = f"L{lecture_num}: {slide_data['title']}"
    title_frame.word_wrap = True

    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Add bullet points text box
    left = Inches(0.8)
    top = Inches(1.5)
    width = Inches(8.4)
    height = Inches(5.5)

    content_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    # Add each bullet point
    for i, bullet in enumerate(slide_data['bullets']):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(64, 64, 64)
        p.space_before = Pt(6)
        p.space_after = Pt(6)

    # Add narration to speaker notes
    if slide_data['narration']:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = slide_data['narration']

def generate_powerpoint():
    """Main function to generate complete PowerPoint presentation."""
    print("🎬 Generating EU AI Act Course PowerPoint...")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Create title slide
    print("📝 Creating title slide...")
    create_title_slide(prs)

    # Process all lecture files
    course_dir = Path('/Users/amarbendou/Documents/Claude/EU AI Act Lab/course')
    lecture_files = sorted(course_dir.glob('Lecture-*.md'))

    total_slides = 0

    for lecture_file in lecture_files:
        print(f"📖 Processing {lecture_file.name}...")

        # Extract lecture data
        lecture_data = extract_lecture_data(lecture_file)

        # Create section divider
        create_section_divider(
            prs,
            lecture_data['lecture_num'],
            lecture_data['lecture_title'],
            lecture_data['duration']
        )

        # Create content slides
        for slide_data in lecture_data['slides']:
            create_content_slide(prs, lecture_data['lecture_num'], slide_data)
            total_slides += 1

    # Save presentation
    output_file = course_dir / 'EU-AI-Act-Course-Complete.pptx'
    prs.save(str(output_file))

    print(f"\n✅ PowerPoint created successfully!")
    print(f"📊 Total slides: {len(prs.slides)} (1 title + 9 section dividers + {total_slides} content slides)")
    print(f"💾 Saved to: {output_file}")
    print(f"📏 File size: {output_file.stat().st_size / 1024:.1f} KB")

if __name__ == '__main__':
    generate_powerpoint()
